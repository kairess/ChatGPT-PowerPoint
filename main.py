import openai
from pptx import Presentation
from pptx.util import Cm
import requests
import io
import json

TOPIC = "한우 채끝살 짜파구리(짜파게티+너구리)를 맛있게 만드는 방법"
NUM_PAGES = 5

openai.api_key = "[YOUR-API-KEY]"

gpt_prompt = [{
    "role": "system",
    "content": (
        "You are a software developer creating an automated PowerPoint generation program."
        " You decide the content that goes into each slide of the PowerPoint."
        " Each slide typically consists of a topic, introduction, main points, conclusion, and references."
        " Follow the rules below:\n"
        f"1. Summarize and extract the key contents from the user's input around {NUM_PAGES} slides.\n"
        "2. Each slide contains 'title', 'content' and 'prompt'.\n"
        "3. 'content' is a bullet point that breaks down the core content into brief, step-by-step chunks.\n"
        "4. All of the slides contain images.\n"
        "5. If the slide contains a image, create a prompt to generate an image using the DALL-E API based on the summarized and extracted content. And save it into 'prompt'.\n"
        "6. Focus on nouns and adjectives and separate them with commas so that 'prompt' is a good visual representation of 'content'.\n"
        "7. Set the above contents as keys named 'title', 'content', and 'prompt'.\n"
        "8. Translate only 'title' and 'content' to Korean. Leave 'prompt' in English without translation.\n"
        "9. If 'prompt' contains Korean, translate it to English.\n"
        "10. Save the results of each slide as a JSON list.\n"
        "11. Output the final output in JSON format.\n"
        "12. Make sure output JSON can be parsed with Python `json.loads()`.\n"
        "13. Must return JSON format only\n\n"
        "Output example:\n\n"
"""```[
  {
    "title": "대한민국 저출산 문제",
    "content": [
        "1. 출산율의 감소",
        "2. 고령화 사회 문제",
        "3. 인구 감소로 인한 경제적 영향",
        "4. 사회적 공헌 부족",
        "5. 가족의 사회적 지원 부족"
    ],
    "prompt": "sad family, old people, small family"
  },
  {
    "title": "캠페인 주제 선택",
    "content": [
        "1. 출산 장려",
        "2. 가족 지원 정책 강화",
        "3. 청소년 성교육 강화",
        "4. 경제적 혜택 제공",
        "5. 사회적 기업의 참여"
    ],
    "prompt": "children, playing in garden, money, family"
  },
  {
    "title": "출산 장려",
    "content": [
        "1. 육아휴직 제도 개선",
        "2. 보육시설 확충",
        "3. 출산 관련 정보 제공",
        "4. 출산 후 복귀를 돕는 프로그램 개발",
        "5. 생명보험 혜택 강화"
    ],
    "prompt": "parents with baby, kindergarten, pregnant, mother"
  },
  {
    "title": "가족 지원 정책 강화",
    "content": [
        "1. 주택 공급 정책 개선",
        "2. 국공립 어린이집 확대",
        "3. 교육 지원 프로그램 신설",
        "4. 가족 세율 혜택 확대",
        "5. 가정 양육지원 서비스 강화"
    ],
    "prompt": "house, kindergarten, school, tax, baby care service"
  }
]```"""
    )
}, {
    "role": "user",
    "content": TOPIC
}]

print("[1] Generating contents...")

gpt_response = openai.ChatCompletion.create(
    model="gpt-4",
    messages=gpt_prompt)

contents = gpt_response["choices"][0]["message"]["content"]

print(contents)

contents = json.loads(contents.replace("`", ""), strict=False)

with open(f"{TOPIC}.json", "w") as f:
    json.dump(contents, f)

print("[2] Generating PowerPoint slides...")

prs = Presentation()

for i, content in enumerate(contents):
    layout = prs.slide_layouts[3] # Left: text, Right: blank

    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = content["title"]

    body_shape = slide.shapes.placeholders[1] # placeholders[0]: Title
    body_shape.text_frame.text = "\n".join(content["content"])

    if "prompt" in content and content["prompt"]:
        print(f"[3-{i+1}] Generating a image for slide #{i+1}...")

        dalle_response = openai.Image.create(
            prompt=content["prompt"],
            size="512x512"
        )

        img_url = dalle_response["data"][0]["url"]

        img_bytes = requests.get(img_url).content
        img_stream = io.BytesIO()
        img_stream.write(img_bytes)

        left = Cm(13)
        top = Cm(4.5)
        height = Cm(11)
        pic = slide.shapes.add_picture(img_stream, left, top, height=height)

print(f"[4] Saving result into disk... {TOPIC}.pptx")
prs.save(f"{TOPIC}.pptx")

print("[5] Done!")
